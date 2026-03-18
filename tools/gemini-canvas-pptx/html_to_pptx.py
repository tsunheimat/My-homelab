import os
import io
import sys
import requests
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Color Palette based on HTML
BG_COLOR = RGBColor(247, 243, 234)     # #f7f3ea - Rice Paper
TITLE_COLOR = RGBColor(126, 25, 27)    # #7e191b - Lacquer Red
SUBTITLE_COLOR = RGBColor(85, 62, 50)  # #553e32 - Dark Brown Subtitle
TEXT_COLOR = RGBColor(51, 37, 30)      # #33251e - Ink Text
ACCENT_COLOR = RGBColor(173, 138, 78)  # #ad8a4e - Gold Highlight

# Box styles
BOX_BG = RGBColor(240, 233, 216)       # #f0e9d8
BOX_BORDER = RGBColor(213, 203, 185)   # #d5cbb9
LIGHT_TEXT = RGBColor(247, 243, 234)

# Fallback Fonts
FONT_SERIF = 'Microsoft JhengHei'
FONT_SANS = 'Microsoft JhengHei'

def apply_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def format_run(run, size, color=TEXT_COLOR, bold=False, font_name=FONT_SANS):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name

def extract_content(container, tf, light_mode=False):
    """Recursively extract content into a text frame."""
    seen = set()
    first = True
    
    # We skip titles that are specifically mapped elsewhere if they are inside container
    for elem in container.find_all(['h3', 'p', 'li', 'div']):
        # Ignore structural divs without text
        if elem.name == 'div' and not any(cls in elem.get('class', []) for cls in ['number', 'number-label']):
            continue
            
        txt = elem.get_text(strip=True)
        if not txt or txt in seen:
            continue
        seen.add(txt)
        
        if first:
            p_run = tf.paragraphs[0]
            first = False
        else:
            p_run = tf.add_paragraph()
            
        p_run.text = "• " + txt if elem.name == 'li' else txt
        
        # Determine formatting based on html tags and classes
        if elem.name == 'h3':
            c = TITLE_COLOR if not light_mode else LIGHT_TEXT
            format_run(p_run.runs[0], 24, c, True, FONT_SERIF)
            p_run.space_after = Pt(10)
        elif 'number' in elem.get('class', []):
            format_run(p_run.runs[0], 80, ACCENT_COLOR, True, FONT_SERIF)
        elif 'number-label' in elem.get('class', []):
            format_run(p_run.runs[0], 20, LIGHT_TEXT, True, FONT_SANS)
            p_run.space_after = Pt(20)
        else:
            c = TEXT_COLOR if not light_mode else LIGHT_TEXT
            format_run(p_run.runs[0], 18, c, False, FONT_SANS)
            p_run.space_after = Pt(10)

def add_image_from_url(slide, img_url, left, top, width=None, height=None):
    try:
        print(f"Downloading image: {img_url}")
        response = requests.get(img_url, timeout=10)
        if response.status_code == 200:
            image_stream = io.BytesIO(response.content)
            slide.shapes.add_picture(image_stream, left, top, width=width, height=height)
    except Exception as e:
        print(f"Failed to load image {img_url}: {e}")

def convert_html_to_pptx(html_file, output_file):
    print(f"Reading HTML file: {html_file}")
    with open(html_file, 'r', encoding='utf-8') as f:
        html_content = f.read()

    soup = BeautifulSoup(html_content, 'html.parser')
    prs = Presentation()
    
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]
    
    slides = soup.find_all('div', class_='slide-container')
    print(f"Found {len(slides)} slides. Converting...")
    
    for slide_div in slides:
        slide = prs.slides.add_slide(blank_layout)
        apply_background(slide)
        
        title_layout = slide_div.find('div', class_='title-layout')
        section_title = slide_div.find('div', class_='section-title-layout')
        qa_layout = slide_div.find('div', class_='qa-layout')
        
        # 1. Main Title Slide
        if title_layout:
            h1 = title_layout.find('h1')
            p = title_layout.find('p')
            
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.333), Inches(1.5), Inches(10.666), Inches(4.5))
            shape.fill.solid()
            shape.fill.fore_color.rgb = BOX_BG
            shape.line.color.rgb = ACCENT_COLOR
            shape.line.width = Pt(2)
            
            tf = shape.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            if h1:
                p_run = tf.paragraphs[0]
                p_run.text = h1.get_text(strip=True, separator='\n')
                p_run.alignment = PP_ALIGN.CENTER
                format_run(p_run.runs[0], 54, TITLE_COLOR, True, FONT_SERIF)
                
            if p:
                p2 = tf.add_paragraph()
                p2.text = p.get_text(strip=True)
                p2.alignment = PP_ALIGN.CENTER
                p2.space_before = Pt(30)
                format_run(p2.runs[0], 24, SUBTITLE_COLOR, False, FONT_SANS)

        # 2. Section Title Slide
        elif section_title:
            h2 = section_title.find('h2')
            p = section_title.find('p')
            txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.33), Inches(2.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            if h2:
                p_run = tf.paragraphs[0]
                p_run.text = h2.get_text(strip=True)
                p_run.alignment = PP_ALIGN.CENTER
                format_run(p_run.runs[0], 48, TITLE_COLOR, True, FONT_SERIF)
                
            if p:
                p2 = tf.add_paragraph()
                p2.text = p.get_text(strip=True)
                p2.alignment = PP_ALIGN.CENTER
                p2.space_before = Pt(20)
                format_run(p2.runs[0], 24, SUBTITLE_COLOR, False, FONT_SANS)

        # 3. QA / Conclusion Layout
        elif qa_layout:
            txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(10.33), Inches(3.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            h2 = qa_layout.find('h2')
            if h2:
                p_run = tf.paragraphs[0]
                p_run.text = h2.get_text(strip=True)
                p_run.alignment = PP_ALIGN.CENTER
                format_run(p_run.runs[0], 48, TITLE_COLOR, True, FONT_SERIF)
                
            p = qa_layout.find('p')
            if p:
                p2 = tf.add_paragraph()
                p2.text = p.get_text(strip=True, separator='\n')
                p2.alignment = PP_ALIGN.CENTER
                p2.space_before = Pt(20)
                format_run(p2.runs[0], 24, SUBTITLE_COLOR, False, FONT_SANS)
                
            contact = qa_layout.find('div', class_='contact-info')
            if contact:
                p3 = tf.add_paragraph()
                p3.text = contact.get_text(strip=True, separator='\n')
                p3.alignment = PP_ALIGN.CENTER
                p3.space_before = Pt(40)
                format_run(p3.runs[0], 22, TITLE_COLOR, True, FONT_SANS)

        # 4. Standard Content Slides
        else:
            slide_title = slide_div.find('h2', class_='slide-title')
            content_area = slide_div.find('div', class_='content-area')
            
            # Setup Header
            if slide_title:
                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8))
                tf = txBox.text_frame
                p_run = tf.paragraphs[0]
                p_run.text = slide_title.get_text(strip=True, separator=' ')
                format_run(p_run.runs[0], 36, TITLE_COLOR, True, FONT_SERIF)

            # Check for layout classes within content area
            is_bleed = 'bleed-image-layout' in slide_div.get('class', [])
            two_column = slide_div.find('div', class_='two-column') if content_area else None
            tiled = slide_div.find('div', class_='tiled-content') if content_area else None
            
            # Bleed Image Layout (Right half is image)
            if is_bleed:
                text_side = slide_div.find('div', class_='bleed-text-side')
                img_tag = slide_div.find('img')
                
                if text_side:
                    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(5.5))
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    extract_content(text_side, tf)
                    
                if img_tag and 'src' in img_tag.attrs:
                    add_image_from_url(slide, img_tag['src'], Inches(6.6), Inches(0), Inches(6.733), Inches(7.5))

            # Tiled 3-Column Content (Icons & Boxes)
            elif tiled:
                tiles = tiled.find_all('div', class_='tile')
                width = 3.8
                spacing = 0.2
                for i, tile in enumerate(tiles):
                    left = Inches(0.5 + i * (width + spacing))
                    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), Inches(width), Inches(5.0))
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = BOX_BG
                    shape.line.color.rgb = BOX_BORDER
                    shape.line.width = Pt(1)
                    
                    tf = shape.text_frame
                    tf.word_wrap = True
                    tf.margin_left = Inches(0.3)
                    tf.margin_right = Inches(0.3)
                    tf.margin_top = Inches(0.4)
                    
                    extract_content(tile, tf)

            # Two Column Content
            elif two_column:
                is_tiled = 'tiled' in two_column.get('class', [])
                is_numbers = 'highlight-numbers-layout' in two_column.get('class', [])
                
                cols = two_column.find_all('div', recursive=False)
                if len(cols) >= 2:
                    for i, col in enumerate(cols):
                        w = 5.8
                        x = 0.5 if i == 0 else 6.8
                        
                        if is_numbers and i == 0:
                            w = 4.5
                            x = 1.0
                        elif is_numbers and i == 1:
                            x = 6.0
                            w = 6.8

                        # Box vs Textbox mapping
                        if is_tiled or (is_numbers and i == 0):
                            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.5), Inches(w), Inches(5.0) if not is_numbers else Inches(4.0))
                            shape.fill.solid()
                            if is_numbers and i == 0:
                                shape.fill.fore_color.rgb = TITLE_COLOR
                                shape.line.fill.background() # Removing border explicitly
                            else:
                                shape.fill.fore_color.rgb = BOX_BG
                                shape.line.color.rgb = BOX_BORDER
                                shape.line.width = Pt(1)
                                
                            tf = shape.text_frame
                            tf.margin_top = Inches(0.4)
                            tf.margin_left = Inches(0.4)
                            tf.margin_right = Inches(0.4)
                            if is_numbers and i == 0:
                                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                                p_run = tf.paragraphs[0]
                                p_run.alignment = PP_ALIGN.CENTER
                            tf.word_wrap = True
                        else:
                            box = slide.shapes.add_textbox(Inches(x), Inches(1.5), Inches(w), Inches(5.5))
                            tf = box.text_frame
                            tf.word_wrap = True
                            
                        # Try to handle images specifically
                        img_tag = col.find('img')
                        if img_tag and 'src' in img_tag.attrs:
                            add_image_from_url(slide, img_tag['src'], Inches(x), Inches(1.5), width=Inches(w))
                        else:
                            extract_content(col, tf, light_mode=(is_numbers and i == 0))
                            
                        if is_numbers and i == 0:
                             for p in tf.paragraphs:
                                  p.alignment = PP_ALIGN.CENTER

            # Fallback simple vertical flow
            elif content_area:
                img_tag = content_area.find('img')
                if img_tag and 'src' in img_tag.attrs:
                    add_image_from_url(slide, img_tag['src'], Inches(8.5), Inches(1.5), width=Inches(4.5))
                    
                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(7.5) if img_tag else Inches(12.33), Inches(5.5))
                tf = txBox.text_frame
                tf.word_wrap = True
                extract_content(content_area, tf)

    prs.save(output_file)
    print(f"Done! Saved presentation to {output_file}")


if __name__ == '__main__':
    if len(sys.argv) > 1:
        html_file = sys.argv[1]
    else:
        html_file = 'test.html'
        
    if not os.path.exists(html_file):
        print(f"Error: Could not find {html_file}")
        sys.exit(1)
        
    output_file = os.path.splitext(html_file)[0] + '.pptx'
    convert_html_to_pptx(html_file, output_file)
